# -*- coding: utf-8 -*-
"""
PPT 生成器模块
使用 python-pptx 生成可编辑的 PPT 文件，复原文字格式
"""

from dataclasses import dataclass, field
from typing import List, Optional, Tuple, Dict, Any
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from .config import PPTConfig
from .ocr_processor import OCRResult, TextFormat


# 常用字体映射（针对 WPS 兼容性）
FONT_NAME_MAPPING = {
    # 英文/无衬线字体 -> WPS 兼容版本
    'Arial': 'Arial',
    'Helvetica': 'Arial',
    'sans-serif': '微软雅黑',
    'San Francisco': '苹方',
    'Segoe UI': '微软雅黑',
    # 中文常用字体
    '微软雅黑': '微软雅黑',
    'Microsoft YaHei': '微软雅黑',
    '宋体': '宋体',
    'SimSun': '宋体',
    '黑体': '黑体',
    'SimHei': '黑体',
    '楷体': '楷体',
    'KaiTi': '楷体',
    '仿宋': '仿宋',
    'FangSong': '仿宋',
    # 默认
    'default': '微软雅黑',
}


@dataclass
class PPTGenerationResult:
    """PPT 生成结果"""
    success: bool = False
    output_path: Optional[str] = None
    slides_count: int = 0
    error_message: Optional[str] = None


class PPTGenerator:
    """PPT 生成器主类"""
    
    def __init__(self, config: Optional[PPTConfig] = None):
        self.config = config or PPTConfig()
        self._presentation = None
    
    def create_presentation(self, width: Optional[float] = None, 
                          height: Optional[float] = None) -> Presentation:
        """创建新的演示文稿
        
        Args:
            width: 幻灯片宽度（英寸），默认使用配置值
            height: 幻灯片高度（英寸），默认使用配置值
        
        Returns:
            Presentation 对象
        """
        self._presentation = Presentation()
        
        # 设置幻灯片尺寸
        slide_width = width or self.config.slide_width
        slide_height = height or self.config.slide_height
        
        # 设置为宽屏 16:9
        self._presentation.slide_width = Inches(slide_width)
        self._presentation.slide_height = Inches(slide_height)
        
        return self._presentation
    
    def add_slide_with_image(self, image: Image.Image,
                            ocr_result: Optional[OCRResult] = None,
                            title: Optional[str] = None) -> int:
        """添加一张幻灯片，包含背景图片和文字
        
        Args:
            image: 背景图片
            ocr_result: OCR 识别结果，用于提取文字格式
            title: 可选的幻灯片标题
        
        Returns:
            幻灯片索引
        """
        if self._presentation is None:
            self.create_presentation()
        
        # 添加幻灯片布局（空白布局）
        slide_layout = self._presentation.slide_layouts[6]  # 空白布局
        slide = self._presentation.slides.add_slide(slide_layout)
        
        # 设置背景图片
        self._set_slide_background(slide, image)
        
        # 添加文字
        if ocr_result:
            self._add_texts_to_slide(slide, ocr_result, image.size)
        
        return len(self._presentation.slides) - 1
    
    def _set_slide_background(self, slide, image: Image.Image):
        """设置幻灯片背景图片"""
        # 将 PIL Image 转换为字节流
        from io import BytesIO
        img_buffer = BytesIO()
        image.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        
        # 添加背景图片
        slide.shapes.add_picture(
            img_buffer,
            Inches(0),
            Inches(0),
            width=Inches(self.config.slide_width),
            height=Inches(self.config.slide_height),
        )
    
    def _add_texts_to_slide(self, slide, ocr_result: OCRResult, 
                           image_size: Tuple[int, int]):
        """添加文字到幻灯片
        
        根据 OCR 结果中的文字格式信息，在对应位置添加可编辑的文字
        """
        img_width, img_height = image_size
        
        for text_format in ocr_result.texts:
            # 计算文字框位置（基于图片坐标比例）
            x1, y1, x2, y2 = text_format.bbox
            
            # 转换为 PPT 坐标（英寸）
            left = (text_format.center_x * self.config.slide_width - 
                   self._estimate_text_width(text_format) / 2)
            top = text_format.center_y * self.config.slide_height
            
            # 限制在幻灯片范围内
            left = max(0, min(left, self.config.slide_width - 0.5))
            top = max(0, min(top, self.config.slide_height - 0.5))
            
            width = Inches(0.5)  # 最小宽度
            height = Inches(0.3)  # 最小高度
            
            # 添加文本框
            textbox = slide.shapes.add_textbox(
                Inches(left),
                Inches(top),
                width,
                height,
            )
            
            # 设置文字内容和格式
            self._format_textbox(textbox, text_format)
    
    def _estimate_text_width(self, text_format: TextFormat) -> float:
        """估算文字框宽度（英寸）"""
        # 基于字数和字体大小估算
        char_count = len(text_format.text)
        char_width = text_format.font_size * 0.5  # 粗略估算
        return max(0.5, char_count * char_width / 72)  # 转换为英寸
    
    def _format_textbox(self, textbox, text_format: TextFormat):
        """格式化文本框"""
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # 设置文字
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text_format.text
        
        # 设置字体
        font = run.font
        font.name = self._map_font_name(text_format.font_name)
        font.size = Pt(text_format.font_size)
        
        # 设置颜色
        r, g, b = text_format.font_color_rgb
        font.color.rgb = RGBColor(r, g, b)
        
        # 设置样式
        font.bold = text_format.bold
        font.italic = text_format.italic
        font.underline = text_format.underline
        
        # 设置对齐方式
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY,
        }
        p.alignment = alignment_map.get(text_format.alignment, PP_ALIGN.LEFT)
        
        # 设置项目符号
        if text_format.bullet:
            self._add_bullet_to_paragraph(p, text_format.bullet)
    
    def _add_bullet_to_paragraph(self, paragraph, bullet: str):
        """添加项目符号到段落"""
        try:
            # 使用 XML 添加项目符号
            pPr = paragraph._p.get_or_add_pPr()
            
            # 创建项目符号定义
            buChar = parse_xml(
                f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet}"/>'
            )
            pPr.append(buChar)
        except Exception:
            # 如果 XML 方法失败，在文字前添加符号
            pass
    
    def _map_font_name(self, font_name: str) -> str:
        """映射字体名称以确保 WPS 兼容性"""
        # 清理字体名称
        font_name = font_name.strip()
        
        # 检查直接映射
        if font_name in FONT_NAME_MAPPING:
            return FONT_NAME_MAPPING[font_name]
        
        # 检查小写
        font_lower = font_name.lower()
        for key, value in FONT_NAME_MAPPING.items():
            if key.lower() == font_lower:
                return value
        
        # 返回默认字体
        return FONT_NAME_MAPPING['default']
    
    def add_slide_with_layout(self, ocr_result: OCRResult,
                              layout_type: str = "title_content") -> int:
        """根据 OCR 分析的布局添加幻灯片
        
        Args:
            ocr_result: OCR 识别结果
            layout_type: 布局类型 "title_content", "two_column", "only_content"
        
        Returns:
            幻灯片索引
        """
        if self._presentation is None:
            self.create_presentation()
        
        slide_layout = self._presentation.slide_layouts[6]  # 空白布局
        slide = self._presentation.slides.add_slide(slide_layout)
        
        page_layout = ocr_result.page_layout
        
        if layout_type == "title_content" or page_layout.get('has_title'):
            # 标题+内容布局
            self._add_title_content_layout(slide, ocr_result)
        
        elif layout_type == "two_column" or page_layout.get('columns', 1) > 1:
            # 双栏布局
            self._add_two_column_layout(slide, ocr_result)
        
        else:
            # 仅内容布局
            self._add_content_only_layout(slide, ocr_result)
        
        return len(self._presentation.slides) - 1
    
    def _add_title_content_layout(self, slide, ocr_result: OCRResult):
        """添加标题+内容布局"""
        title_texts = ocr_result.title_texts
        body_texts = ocr_result.body_texts
        
        # 添加标题
        for title in title_texts:
            self._add_text_box(slide, title, is_title=True)
        
        # 添加正文
        for body in body_texts:
            self._add_text_box(slide, body, is_title=False)
    
    def _add_two_column_layout(self, slide, ocr_result: OCRResult):
        """添加双栏布局"""
        texts = sorted(ocr_result.texts, key=lambda t: t.center_x)
        mid_x = 0.5  # 中间分割线
        
        for text in texts:
            is_left = text.center_x < mid_x
            self._add_text_box(slide, text, column=('left' if is_left else 'right'))
    
    def _add_content_only_layout(self, slide, ocr_result: OCRResult):
        """添加仅内容布局"""
        for text in ocr_result.texts:
            self._add_text_box(slide, text)
    
    def _add_text_box(self, slide, text_format: TextFormat,
                      is_title: bool = False, column: str = None):
        """添加文本框"""
        # 计算位置
        x = text_format.center_x * self.config.slide_width
        y = text_format.center_y * self.config.slide_height
        
        # 根据列调整
        if column == 'left':
            x = x * 0.45  # 左栏
        elif column == 'right':
            x = 0.5 + x * 0.45  # 右栏
        
        # 设置宽度
        width = self.config.slide_width - Inches(1)
        if column:
            width = width / 2 - Inches(0.25)
        
        # 添加文本框
        textbox = slide.shapes.add_textbox(
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(0.5),
        )
        
        self._format_textbox(textbox, text_format)
    
    def add_slide_from_image_with_text(self, image: Image.Image,
                                       text_formats: List[TextFormat]) -> int:
        """从图片和文字格式列表添加幻灯片
        
        Args:
            image: 背景图片
            text_formats: 文字格式列表
        
        Returns:
            幻灯片索引
        """
        if self._presentation is None:
            self.create_presentation()
        
        slide_layout = self._presentation.slide_layouts[6]
        slide = self._presentation.slides.add_slide(slide_layout)
        
        # 设置背景
        self._set_slide_background(slide, image)
        
        # 添加文字
        for tf in text_formats:
            self._add_text_box(slide, tf)
        
        return len(self._presentation.slides) - 1
    
    def save(self, output_path: str) -> PPTGenerationResult:
        """保存演示文稿
        
        Args:
            output_path: 输出文件路径
        
        Returns:
            PPTGenerationResult: 生成结果
        """
        if self._presentation is None:
            return PPTGenerationResult(
                success=False,
                error_message="没有可保存的演示文稿",
            )
        
        try:
            # 确保输出目录存在
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            
            # 保存文件
            self._presentation.save(output_path)
            
            return PPTGenerationResult(
                success=True,
                output_path=output_path,
                slides_count=len(self._presentation.slides),
            )
        
        except Exception as e:
            return PPTGenerationResult(
                success=False,
                error_message=f"保存 PPT 时出错: {str(e)}",
            )
    
    def get_presentation(self) -> Optional[Presentation]:
        """获取当前的演示文稿对象"""
        return self._presentation


class PPTGeneratorPipeline:
    """PPT 生成流水线 - 整合 OCR 和图像修复结果"""
    
    def __init__(self, config: Optional[PPTConfig] = None):
        self.ppt_config = config or PPTConfig()
        self.generator = PPTGenerator(self.ppt_config)
    
    def create_from_ocr_and_inpaint(self,
                                   image: Image.Image,
                                   ocr_result: OCRResult,
                                   cleaned_image: Optional[Image.Image] = None,
                                   output_path: Optional[str] = None) -> PPTGenerationResult:
        """从 OCR 结果和修复图片创建 PPT
        
        Args:
            image: 原始图片
            ocr_result: OCR 识别结果
            cleaned_image: 修复后的图片（可选，默认使用原始图片）
            output_path: 输出文件路径
        
        Returns:
            PPTGenerationResult: 生成结果
        """
        # 使用修复后的图片作为背景
        bg_image = cleaned_image or image
        
        # 创建演示文稿
        self.generator.create_presentation()
        
        # 添加幻灯片
        self.generator.add_slide_with_image(bg_image, ocr_result)
        
        # 保存
        if output_path:
            return self.generator.save(output_path)
        
        return PPTGenerationResult(
            success=True,
            slides_count=1,
        )
    
    def create_multi_slide(self,
                          slides_data: List[Dict[str, Any]],
                          output_path: str) -> PPTGenerationResult:
        """创建多页 PPT
        
        Args:
            slides_data: 每页的数据列表，每个元素包含:
                - 'image': 背景图片
                - 'ocr_result': OCR 结果（可选）
                - 'text_formats': 文字格式列表（可选）
            output_path: 输出文件路径
        
        Returns:
            PPTGenerationResult: 生成结果
        """
        self.generator.create_presentation()
        
        for slide_data in slides_data:
            image = slide_data.get('image')
            ocr_result = slide_data.get('ocr_result')
            text_formats = slide_data.get('text_formats')
            
            if text_formats and image:
                self.generator.add_slide_from_image_with_text(image, text_formats)
            elif ocr_result and image:
                self.generator.add_slide_with_image(image, ocr_result)
            elif image:
                self.generator.add_slide_with_image(image)
        
        return self.generator.save(output_path)
